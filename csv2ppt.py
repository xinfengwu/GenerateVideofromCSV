
from pptx import Presentation
import csv
import copy
from pptx.util import Pt
import aspose.slides as slides
import aspose.pydrawing as drawing

        
# 定义一个函数，将数据填入幻灯片
def fill_slide(slide, data):

    slide_title = find_shape_by_name(slide.shapes,'index')
    add_text(slide_title,data['index'])
    slide_title = find_shape_by_name(slide.shapes,'hiragana')
    add_text(slide_title,data['phonetic symbol'])
    slide_title = find_shape_by_name(slide.shapes,'kanji')
    add_text(slide_title,data['word'])
         
# 复制幻灯片
def duplicate_slide(prs, index):
    slide_to_copy= prs.slides[index]
    # Create a new slide object
    new_slide = prs.slides.add_slide(slide_to_copy.slide_layout)
    # Copy content from the original slide to the new slide
    for shape in slide_to_copy.shapes:
        el = shape.element
        newel = copy.deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    return new_slide


# 查找文本框    
def find_shape_by_name(shapes, name):
    for shape in shapes:
        if shape.name == name:
            return shape
    return None


# 修改文字
def add_text(shape, text, alignment=None):
    
    if alignment:
        shape.vertical_anchor = alignment

    tf = shape.text_frame
    tf.clear()
    run = tf.paragraphs[0].add_run()
    run.text = text if text else ''
    font=run.font
    font.name = 'Calibri'
    font.size = Pt(58)
    font.bold = True
    
    
# 读取 CSV 文件
csv_file = 'data.csv'  # 请替换为你的 CSV 文件路径

input_file = csv.DictReader(open(csv_file))

# 加载 PowerPoint 文件
ppt_template = 'template.pptx'  # 请替换为你的 PowerPoint 模板文件路径
presentation = Presentation(ppt_template)

# 遍历 CSV 文件中的每条记录，创建对应的幻灯片
for row in input_file:
    # 从模板创建一张新的幻灯片
    new_slide = duplicate_slide(presentation, 0)
    # 填充文本
    fill_slide(new_slide, row)

# 保存修改后的幻灯片
output_ppt = 'output.pptx'  # 输出文件路径
presentation.save(output_ppt)
print(f'幻灯片已保存为 {output_ppt}')

# 将PPTX导出为图片
pres = slides.Presentation("output.pptx")

for sld in pres.slides:
    bmp = sld.get_thumbnail(1, 1)
    bmp.save("Slide_{num}.jpg".format(num=str(sld.slide_number)), drawing.imaging.ImageFormat.jpeg)

print(f'幻灯片已导出为图片 ')
