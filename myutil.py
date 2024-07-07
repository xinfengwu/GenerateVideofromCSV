import os
import gtts
from pptx import Presentation
import csv
import copy
import subprocess
import tempfile
from pdf2image import convert_from_path
from pdf2image.exceptions import (
    PDFInfoNotInstalledError,
    PDFPageCountError,
    PDFSyntaxError
)

import subprocess
from PIL import Image
from pptx.util import Pt
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import shutil


# 创建文件夹
def create_folder(folder_path):
    # 检查目录是否存在
    if not os.path.exists(folder_path):
        # 如果不存在，创建目录
        os.makedirs(folder_path)
        print(f"文件夹 '{folder_path}' 已创建！")
    else:
        print(f"文件夹 '{folder_path}' 已存在！")
        empty_folder(folder_path)
        print(f"文件夹 '{folder_path}' 已清空！")
    return folder_path
    

# 删除空文件
def empty_folder(folder_path):
    # 遍历文件夹中的所有文件
    for file_name in os.listdir(folder_path):
        file_path = os.path.join(folder_path, file_name)

        # 如果是空文件，则删除
        if os.path.isfile(file_path) and os.path.getsize(file_path) == 0: # 空文件
            os.remove(file_path)
        # 如果是文件夹，则递归调用 empty_folder 函数清空该文件夹
        elif os.path.isdir(file_path):
            empty_folder(file_path)
            
            
# 读取csv文件
def read_csv_data(file_path):
    data = []
    with open(file_path, newline='') as csvfile:
        reader = csv.DictReader(csvfile)
        for row in reader:
            data.append(row)
    return data     


def create_ppt_with_csv(data, ppt_template, bg_img_path, output_ppt):
    prs = Presentation(ppt_template)

    # 遍历 CSV 文件中的每条记录，创建对应的幻灯片
    for row in data:
        # 从模板创建一张新的幻灯片
        new_slide = duplicate_slide(prs, 0, bg_img_path)
        # 填充文本
        fill_slide(new_slide, row)

    # 删除模版幻灯片
    delete_slide(prs,0)
    # 保存修改后的幻灯片
    prs.save(output_ppt)
    # print(f'幻灯片已保存为 {output_ppt}')
    return prs

   
# 复制幻灯片
def duplicate_slide(prs, index, bg_img_path):
    slide_to_copy = prs.slides[index]

    # Create a new slide object
    new_slide = prs.slides.add_slide(slide_to_copy.slide_layout)
    
    # 设置背景图片
    if bg_img_path:
        set_background_image(prs, new_slide, bg_img_path)
    
    # create images dict
    imgDict = {}
    
    # 复制各种形状
    for shape in slide_to_copy.shapes:
        # 使用 shape_type 属性判断形状类型
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # print("This is a picture.")
            # save image
            with open(shape.name+'.jpg', 'wb') as f:
                f.write(shape.image.blob)
            # add image to dict
            imgDict[shape.name+'.jpg'] = [shape.left, shape.top, shape.width, shape.height]
        elif shape.shape_type == 1: # MSO_SHAPE_TYPE.AUTO_SHAPE:
            # print("This is a AUTO_SHAPE.")
            el = shape.element
            newel = copy.deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
        else:
            print(f"This is an unhandled shape type: {shape.shape_type}")
            
    
    # add pictures
    for k, v in imgDict.items():
        new_slide.shapes.add_picture(k, v[0], v[1], v[2], v[3])
        os.remove(k)

    return new_slide
    
    
# 将数据填入幻灯片
def fill_slide(slide, row):
    slide_title = find_shape_by_name(slide.shapes,'index')
    add_text(slide_title,row['index'])
    slide_title = find_shape_by_name(slide.shapes,'hiragana')
    add_text(slide_title,row['hiragana'])
    slide_title = find_shape_by_name(slide.shapes,'kanji')
    add_text(slide_title,row['kanji'])


# 查找文本框
def find_shape_by_name(shapes, name):
    for shape in shapes:
        # print(shape.shape_id)
        if shape.name == name:
            return shape
            

# 修改文字
def add_text(shape, text,):
    tf = shape.text_frame
    p = tf.paragraphs[0].runs
    p[0].text = text if text else ''

   
# 设置幻灯片背景
def set_background_image(prs, new_slide, bg_img_path):
    left = top = Inches(0)
    pic = new_slide.shapes.add_picture(bg_img_path, left, top, width=prs.slide_width, height=prs.slide_height)

    # This moves it to the background
    new_slide.shapes._spTree.remove(pic._element)
    new_slide.shapes._spTree.insert(2, pic._element)
   
   
# 删除幻灯片
def delete_slide(prs, slide_index):
    xml_slides = prs.slides._sldIdLst  # Access the slide list XML elements
    if slide_index < 0 or slide_index >= len(xml_slides):
        print(f"Slide index {slide_index} out of range.")
        return
    slide_id = xml_slides[slide_index].rId
    prs.slides._sldIdLst.remove(xml_slides[slide_index])  # Remove slide from XML elements
    prs.part.drop_rel(slide_id)  # Remove slide from related parts
    
    
# 将ppt转化成pdf
# pdf页面和对应幻灯片相比会有一点变形
def ppt_to_pdf_by_soffice(ppt_file, pdf_folder, pdf_file):
    # kill all the soffice
    kill_soffice_cmd = "killall soffice.bin"
    soffice_cmd = "soffice --headless --convert-to pdf " + ppt_file + " --outdir " + pdf_folder
    subprocess.call(soffice_cmd, shell=True)

    return pdf_file
    

def ppt_to_pdf_by_unoconv(ppt_file, body_pdf):
    # subprocess.call(['unoconv', '-f', 'pdf', '-o', pdf_folder, ppt_file])
    unoconv_cmd = "unoconv -f pdf -o " + body_pdf + " " + ppt_file
    subprocess.call(unoconv_cmd, shell=True)

# 将pdf转化成jpg
def pdf_to_img(pdf_file, imgs_folder):
    pages = convert_from_path(pdf_file, first_page=0)
    for img in pages:
        img_path = imgs_folder + '/' + str(pages.index(img)+1) + ".jpg"
        img.save(img_path , quality=100)
    # print(f'PDF已导出为图片 ')


def compare_mp3_and_img_files(mp3_folder, img_folder):
    # 筛选出 mp3_folder 中所有的 mp3 文件和 img_folder 中所有的 img 文件
    mp3_files = filter_files_by_extension(mp3_folder, '.mp3')
    img_files = filter_files_by_extension(img_folder, '.jpg')
    
    # 去掉后缀,将文件名列表转换为集合，以便进行快速比较
    mp3_set = {file.split('.')[0] for file in mp3_files}
    img_set = {file.split('.')[0] for file in img_files}
    
    # 找出两个集合中共同的文件名
    common_files = mp3_set.intersection(img_set)
    
    # 找出只在 mp3 文件夹中的文件名
    unique_to_mp3 = mp3_set - img_set
    
    # 找出只在 img 文件夹中的文件名
    unique_to_img = img_set - mp3_set
    # 打印结果
    # print(f'共同的文件名: {common_files}')
    # print(f'只在 mp3 文件夹中的文件名: {unique_to_mp3}')
    # print(f'只在 img 文件夹中的文件名: {unique_to_img}')

    return common_files, unique_to_mp3, unique_to_img
    

def filter_files_by_extension(folder, extension):
    # 获取指定文件夹中特定后缀名的文件列表
    file_list = [file for file in os.listdir(folder) if file.endswith(extension)]
    return file_list
    

# 按行生成语音
def text_to_mp3(keyword, language, output_file):
    if not os.path.exists(output_file):
        try:
        # print(gtts.lang.tts_langs()) 输出支持的语言
            tts = gtts.gTTS(keyword, lang=language)  ##  request google to get synthesis
            tts.save(output_file)  #  save audio
        except Exception as e:
            print(f"转换 '{row['kanji'].strip()}' 出现错误：{e}")


def set_video_resolution(slide_width, slide_height):
    # 常量
    EMU_PER_INCH = 914400
    PIXELS_PER_INCH = 96

    # 获取幻灯片宽度和高度（EMU）
    slide_width_emu = slide_width
    slide_height_emu = slide_height

    # 将 EMU 转换为英寸
    slide_width_inches = slide_width_emu / EMU_PER_INCH
    slide_height_inches = slide_height_emu / EMU_PER_INCH

    # 将英寸转换为像素
    slide_width_pixels = int(slide_width_inches * PIXELS_PER_INCH)
    slide_height_pixels = int(slide_height_inches * PIXELS_PER_INCH)

    # 视频解析度要为2的倍数
    if slide_width_pixels % 2 != 0:
        slide_width_pixels += 1
    if slide_height_pixels % 2 != 0:
        slide_height_pixels += 1
        
    print(f"Slide width: {slide_width_emu} EMU, {slide_width_pixels:.2f} pixels")
    print(f"Slide height: {slide_height_emu} EMU, {slide_height_pixels:.2f} pixels")
    
    resolution = [slide_width_pixels,slide_height_pixels]
    return resolution   
    
    
# img + mp3 ---> video
def img_mp3_to_mp4(mp3_file, img_file, size, mp4_file):
    # print("img_mp3_to_mp4")
    ffmpeg_cmd = (
        f"ffmpeg -loglevel error  -i {mp3_file} "
        f" -loop 1 -i {img_file}"
        f" -vcodec libx264 -pix_fmt yuv420p -shortest -y "
        f" -vf 'scale={size[0]}:{size[1]}' {mp4_file}"
    )
    # print(ffmpeg_cmd)
    subprocess.call(ffmpeg_cmd, shell=True)
    
    
# img ---> video
def image_to_video(image_path, size, duration, output_file):
    # FFmpeg 命令
    ffmpeg_command = (
        f"ffmpeg -loglevel error -loop 1 -t {duration} -i {image_path}"
        f" -f lavfi -i anullsrc=cl=mono:r=24000 -shortest -y "
        f" -vf 'scale={size[0]}:{size[1]}' {output_file}"
        
    )
    
    # 运行 FFmpeg 命令
    subprocess.run(ffmpeg_command, shell=True, check=True)   
   
    
# 创建mp4_filelist.txt
def create_mp4_filelist(folder_path):
    # 获取文件夹中所有的 MP4 文件,
    mp4_filenames_set = get_mp4_filenames(folder_path)
    
    # 排序
    sorted_mp4_filenames = sort_filenames_as_integers(mp4_filenames_set)
    
    # 组成文件完整路径
    mp4_files = add_path_to_filenames(folder_path, sorted_mp4_filenames)
    
    return mp4_files


# 获取文件夹下所有文件的文件名和扩展名，并保存到set中
def get_mp4_filenames(folder_path):
    filenames_set = set()
    
    # 遍历文件夹中的所有文件
    for root, dirs, files in os.walk(folder_path):
        for file in files:
            filename, extension = os.path.splitext(file)
            if extension == ".mp4":
                filenames_set.add(filename)
    
    return filenames_set
    
    
# 将set里元素转换成整型后排序
def sort_filenames_as_integers(filenames_set):
    # 将文件名转换为整型并排序
    try:
        filenames_list = sorted(int(filename) for filename in filenames_set)
    except ValueError:
        raise ValueError("Some filenames are not integers and cannot be converted.")
    
    # 转换回字符串（如果需要保持字符串格式）
    sorted_filenames = [str(filename) for filename in filenames_list]
    
    return sorted_filenames
    
    
def add_path_to_filenames(folder_path, filenames):
    full_paths = []
    for filename in filenames:
        full_path = os.path.join(folder_path, filename + '.mp4')
        full_paths.append(full_path)
    return full_paths    
    
    
# 将列表中的每个元素重复 n 次
def repeat_elements(input_list, n):
    # 使用列表推导式将每个元素重复 n 次
    return [item for item in input_list for _ in range(n)]
    
    
# 将列表中的每个元素重复n次，并加上间隔视频
def add_elements_with_silence(input_list, silence_file):
    result = []
    for item in input_list:
        result.append(item)
        result.append(silence_file)
    
    return result
    
    
# 将多个短视频合成为一个视频
def concatenate_mp4_files(input_txt, output_file):
    
    # 构建 FFmpeg 命令
    ffmpeg_command = "ffmpeg -loglevel error -f concat -segment_time_metadata 1 -safe 0 -i " \
        + input_txt + \
        " -vf select=concatdec_select -af aselect=concatdec_select,aresample=async=1 -y "\
        + output_file
    
    # 运行FFmpeg命令合并视频并覆盖已有的output.mp4
    subprocess.call(ffmpeg_command, shell=True)
    print("Videos have been merged into output.mp4")    
            
        
# 删除单个文件
def delete_file(file_path):
    if os.path.exists(file_path):
        os.remove(file_path)
        print(f"{file_path} has been deleted")
    else:
        print(f"{file_path} does not exist")


# 删除文件夹及其内容
def delete_folder(folder_path):
    if os.path.exists(folder_path):
        shutil.rmtree(folder_path)
        print(f"{folder_path} and its contents have been deleted")
    else:
        print(f"{folder_path} does not exist")
   

# 创建空白静音视频
def create_silent_video(duration, size, output_file):
    # 创建一个白色的静音视频
    ffmpeg_command = (
        f"ffmpeg -loglevel error -f lavfi -i color=c=white:s={size[0]}x{size[1]}:d={duration} "
        f"-f lavfi -t {duration} -i anullsrc=r=24000:cl=mono "
        f" -c:v libx264 -c:a aac -y {output_file}"
    )
    # -i input.mp4：指定输入视频文件。
    # -f lavfi -t 3 -i anullsrc=r=44100:cl=stereo：
    # -f lavfi：指定输入格式为Libavfilter。
    # -t 3：将输入流的持续时间限制为3秒。
    # -i anullsrc=r=44100:cl=stereo：使用anullsrc滤镜生成一个静音音频流，采样率为44100 Hz，立体声。
    # -c:v copy：直接复制视频流，不重新编码。
    # -c:a aac：将音频流编码为AAC格式。
    subprocess.run(ffmpeg_command, shell=True, check=True)  
       

def create_silent_mp3(duration, output_file):
    # Construct the ffmpeg command
    command = [
        'ffmpeg', 
        '-loglevel', 'error',
        '-f', 'lavfi', 
        '-i', 'anullsrc=r=44100:cl=stereo', 
        '-t', str(duration), 
        '-y', # 覆盖原有文件
        output_file
    ]

    # Run the command
    subprocess.run(command, check=True)
    print("Done! Create silent mp3")
    return output_file
