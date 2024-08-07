import os
import gtts
from pptx import Presentation
# pip install python-pptx

import csv
import copy
import subprocess
import tempfile
from pdf2image import convert_from_path
from pdf2image.exceptions import (
    PDFInfoNotInstalledError,
    PDFPageCountError,
    PDFSyntaxError
)
# pip install pdf2image

import subprocess
from PIL import Image
from pptx.util import Pt
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import shutil


# 创建文件夹
def create_folder(folder_path):
    # 检查目录是否存在
    if not os.path.exists(folder_path):
        # 如果不存在，创建目录
        os.makedirs(folder_path)
        print(f"文件夹 '{folder_path}' 已创建！")
    else:
        print(f"文件夹 '{folder_path}' 已存在！")
        empty_folder(folder_path)
        print(f"文件夹 '{folder_path}' 里的空文件已删除！")
    return folder_path
    

# 删除空文件
def empty_folder(folder_path):
    # 遍历文件夹中的所有文件
    for file_name in os.listdir(folder_path):
        file_path = os.path.join(folder_path, file_name)

        # 如果是空文件，则删除
        if os.path.isfile(file_path) and os.path.getsize(file_path) == 0: # 空文件
            os.remove(file_path)
        # 如果是文件夹，则递归调用 empty_folder 函数清空该文件夹
        elif os.path.isdir(file_path):
            empty_folder(file_path)
            
            
# 读取csv文件
def read_csv_data(file_path):
    data = []
    with open(file_path, newline='') as csvfile:
        reader = csv.DictReader(csvfile)
        for row in reader:
            data.append(row)
    return data     


def create_ppt_with_csv(data, src_ppt_prs, slide_index, bg_img_path, output_ppt):
    """
    slide_index:
        1: 第1张幻灯片
        2: 第2张幻灯片
        3: 第3张幻灯片
        ...
    """
    
    # prs = Presentation(ppt_template)
    # 创建一个新的PPTX文件
    new_prs = Presentation()
    # 设置新PPTX的方向和源PPTX一致
    new_prs.slide_width = src_ppt_prs.slide_width
    new_prs.slide_height = src_ppt_prs.slide_height

    # 遍历 CSV 文件中的每条记录，创建对应的幻灯片
    for row in data:
        # 从模板创建一张新的幻灯片
        new_slide = duplicate_slide(src_ppt_prs, slide_index, new_prs, bg_img_path)
        # 填充文本
        if slide_index == 1:
            fill_slide1(new_slide, row)
        elif slide_index == 2:
            fill_slide2(new_slide, row)
        elif slide_index == 3:
            fill_slide3(new_slide, row)
        elif slide_index == 4:
            fill_slide4(new_slide, row)
        elif slide_index == 5:
            fill_slide5(new_slide, row)
        else:
            print("ppt_type 值异常")

    # 删除模版幻灯片
    # delete_slide(prs,0)
    # 保存修改后的幻灯片
    new_prs.save(output_ppt)
    # print(f'幻灯片已保存为 {output_ppt}')
    return new_prs

   
# 复制幻灯片
def duplicate_slide(src_ppt_prs, index, new_ppt_prs, bg_img_path=""):
    slide_to_copy = src_ppt_prs.slides[int(index)-1]
    # print(slide_to_copy.slide_id)
    # print(src_ppt_prs.slides[index].name)

    # Create a new slide object
    slide_layout = new_ppt_prs.slide_layouts[5]  # 使用一个空白布局
    # new_slide = new_ppt_prs.slides.add_slide(slide_to_copy.slide_layout)
    new_slide = new_ppt_prs.slides.add_slide(slide_layout)
    
    # 设置背景图片
    if os.path.isfile(bg_img_path) and is_image_file(bg_img_path):
        set_background_image(new_ppt_prs, new_slide, bg_img_path)
    
    # create images dict
    imgDict = {}
    
    # 复制各种形状
    for shape in slide_to_copy.shapes:
        # 使用 shape_type 属性判断形状类型
        """
            1: MSO_SHAPE_TYPE.AUTO_SHAPE
            6: GROUP
            13: MSO_SHAPE_TYPE.PICTURE
            17: MSO_SHAPE_TYPE.TEXT_BOX
            
        """
        #print(shape.shape_type)
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # print("This is a picture.")
            # save image
            with open(shape.name+'.jpg', 'wb') as f:
                f.write(shape.image.blob)
            # add image to dict
            imgDict[shape.name+'.jpg'] = [shape.left, shape.top, shape.width, shape.height]
        elif shape.shape_type == 1 or shape.shape_type == 6 or shape.shape_type == 17:
            # print("This is a AUTO_SHAPE.")
            el = shape.element
            newel = copy.deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
        else:
            print(f"This is an unhandled shape type: {shape.shape_type}")
            
    
    # add pictures
    for k, v in imgDict.items():
        new_slide.shapes.add_picture(k, v[0], v[1], v[2], v[3])
        os.remove(k)

    return new_slide
    
    
def is_image_file(file_path):
    # 定义常见的图像文件扩展名
    image_extensions = ('.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.tif', '.webp')
    
    # 检查文件扩展名是否在定义的图像扩展名中
    return file_path.lower().endswith(image_extensions)
    
    
# 将数据填入幻灯片
def fill_slide1(slide, row):
    slide_title = find_shape_by_name(slide.shapes,'Slide_1_index')
    add_text(slide_title,row['序号'])
    slide_title = find_shape_by_name(slide.shapes,'Slide_1_hiragana')
    add_text(slide_title,row['平假名注音'])
    slide_title = find_shape_by_name(slide.shapes,'Slide_1_kanji')
    add_text(slide_title,row['日本語'])
    slide_title = find_shape_by_name(slide.shapes,'Slide_1_english')
    add_text(slide_title,row['英文'])
    slide_title = find_shape_by_name(slide.shapes,'Slide_1_chinese')
    add_text(slide_title,row['中文'])
    
def fill_slide2(slide, row):
    slide_title = find_shape_by_name(slide.shapes,'Slide_2_Lesson')
    add_text(slide_title,row['lesson_name'])
    
def fill_slide3(slide, row):
    slide_title = find_shape_by_name(slide.shapes,'Slide_3_Lesson_title')
    add_text(slide_title,row['lesson_name'])
    
def fill_slide4(slide, row):
    slide_title = find_shape_by_name(slide.shapes,'Slide_4_Book_title')
    add_text(slide_title,row['lesson_name'])

def fill_slide5(slide, row):
    slide_title = find_shape_by_name(slide.shapes,'Slide_5_index')
    add_text(slide_title,row['序号'])
    slide_title = find_shape_by_name(slide.shapes,'Slide_5_A_context')
    add_text(slide_title,row['問句'])
    slide_title = find_shape_by_name(slide.shapes,'Slide_5_B_context')
    add_text(slide_title,row['答句'])
    
    slide_title = find_shape_by_name(slide.shapes,'Slide_5_A_context_hiragana')
    add_text(slide_title,row['問句-平仮名'])
    slide_title = find_shape_by_name(slide.shapes,'Slide_5_B_context_hiragana')
    add_text(slide_title,row['答句-平仮名'])
    slide_title = find_shape_by_name(slide.shapes,'Slide_5_A_context_en')
    add_text(slide_title,row['問句-英语'])
    slide_title = find_shape_by_name(slide.shapes,'Slide_5_A_context_cn')
    add_text(slide_title,row['問句-中文'])
    slide_title = find_shape_by_name(slide.shapes,'Slide_5_B_context_en')
    add_text(slide_title,row['答句-英语'])
    slide_title = find_shape_by_name(slide.shapes,'Slide_5_B_context_cn')
    add_text(slide_title,row['答句-中文'])


# 查找文本框
def find_shape_by_name(shapes, name):
    for shape in shapes:
        # print(shape.shape_id)
        if shape.name == name:
            return shape
            

# 修改文字
def add_text(shape, text,):
    tf = shape.text_frame
    p = tf.paragraphs[0].runs
    p[0].text = ""
    p[0].text = text if text else ''

   
# 设置幻灯片背景
def set_background_image(prs, new_slide, bg_img_path):
    left = top = Inches(0)
    pic = new_slide.shapes.add_picture(bg_img_path, left, top, width=prs.slide_width, height=prs.slide_height)

    # This moves it to the background
    new_slide.shapes._spTree.remove(pic._element)
    new_slide.shapes._spTree.insert(2, pic._element)
   
   
# 删除幻灯片
def delete_slide(prs, slide_index):
    xml_slides = prs.slides._sldIdLst  # Access the slide list XML elements
    if slide_index < 0 or slide_index >= len(xml_slides):
        print(f"Slide index {slide_index} out of range.")
        return
    slide_id = xml_slides[slide_index].rId
    prs.slides._sldIdLst.remove(xml_slides[slide_index])  # Remove slide from XML elements
    prs.part.drop_rel(slide_id)  # Remove slide from related parts
    
    
# 将ppt转化成pdf
# pdf页面和对应幻灯片相比会有一点变形
def ppt_to_pdf_by_soffice(ppt_file, pdf_folder, pdf_file):
    # kill all the soffice
    kill_soffice_cmd = "killall soffice.bin"
    soffice_cmd = "soffice --headless --convert-to pdf " + ppt_file + " --outdir " + pdf_folder
    subprocess.call(soffice_cmd, shell=True)

    return pdf_file
    

def ppt_to_pdf_by_unoconv(ppt_file, output_pdf):
    # subprocess.call(['unoconv', '-f', 'pdf', '-o', pdf_folder, ppt_file])
    unoconv_cmd = "unoconv -f pdf -o " + output_pdf + " " + ppt_file
    subprocess.call(unoconv_cmd, shell=True)

# 将pdf转化成jpg
def pdf_to_img(pdf_file, imgs_folder):
    pages = convert_from_path(pdf_file, first_page=0)
    for img in pages:
        # print(pages.index(img))
        img_path = imgs_folder + '/' + str(pages.index(img)+1) + ".jpg"
        img.save(img_path , quality=100)
    # print(f'PDF已导出为图片 ')


def compare_mp3_and_img_files(mp3_folder, img_folder):
    # 筛选出 mp3_folder 中所有的 mp3 文件和 img_folder 中所有的 img 文件
    mp3_files = filter_files_by_extension(mp3_folder, '.mp3')
    img_files = filter_files_by_extension(img_folder, '.jpg')
    
    # 去掉后缀,将文件名列表转换为集合，以便进行快速比较
    mp3_set = {file.split('.')[0] for file in mp3_files}
    img_set = {file.split('.')[0] for file in img_files}
    
    # 找出两个集合中共同的文件名
    common_files = mp3_set.intersection(img_set)
    
    # 找出只在 mp3 文件夹中的文件名
    unique_to_mp3 = mp3_set - img_set
    
    # 找出只在 img 文件夹中的文件名
    unique_to_img = img_set - mp3_set
    # 打印结果
    # print(f'共同的文件名: {common_files}')
    # print(f'只在 mp3 文件夹中的文件名: {unique_to_mp3}')
    # print(f'只在 img 文件夹中的文件名: {unique_to_img}')

    return common_files, unique_to_mp3, unique_to_img
    

def filter_files_by_extension(folder, extension):
    # 获取指定文件夹中特定后缀名的文件列表
    file_list = [file for file in os.listdir(folder) if file.endswith(extension)]
    return file_list
    

# 按行生成语音
def text_to_mp3(keyword, language, output_file):
    # mp3文件不存在或是空文件
    if (not os.path.isfile(output_file)) or os.path.getsize(output_file) == 0:
        try:
            # print(gtts.lang.tts_langs()) 输出支持的语言
            # Create a gTTS object
            tts = gtts.gTTS(keyword, lang=language, slow=False)  #  request google to get synthesis
            tts.save(output_file)  #  save audio
        except Exception as e:
            print(f"转换 '{keyword}' 出现错误：{e}")


def set_video_resolution(slide_width, slide_height):
    # 常量
    EMU_PER_INCH = 914400
    PIXELS_PER_INCH = 96

    # 获取幻灯片宽度和高度（EMU）
    slide_width_emu = slide_width
    slide_height_emu = slide_height

    # 将 EMU 转换为英寸
    slide_width_inches = slide_width_emu / EMU_PER_INCH
    slide_height_inches = slide_height_emu / EMU_PER_INCH

    # 将英寸转换为像素
    slide_width_pixels = int(slide_width_inches * PIXELS_PER_INCH)
    slide_height_pixels = int(slide_height_inches * PIXELS_PER_INCH)

    # 视频解析度要为2的倍数
    if slide_width_pixels % 2 != 0:
        slide_width_pixels += 1
    if slide_height_pixels % 2 != 0:
        slide_height_pixels += 1
        
    print(f"Slide width: {slide_width_emu} EMU, {slide_width_pixels:.2f} pixels")
    print(f"Slide height: {slide_height_emu} EMU, {slide_height_pixels:.2f} pixels")
    
    resolution = [slide_width_pixels,slide_height_pixels]
    return resolution   
    
    
# img + mp3 ---> video
def img_mp3_to_mp4(mp3_file, img_file, size, mp4_file):
    # print("img_mp3_to_mp4")
    ffmpeg_cmd = (
        f"ffmpeg -loglevel error  -i {mp3_file} "
        f" -loop 1 -i {img_file}"
        f" -vcodec libx264 -pix_fmt yuv420p -shortest -y "
        f" -vf 'scale={size[0]}:{size[1]}' {mp4_file}"
    )
    # print(ffmpeg_cmd)
    subprocess.call(ffmpeg_cmd, shell=True)
    
    
# img ---> video
def image_to_video(image_path, size, duration, output_file):
    # FFmpeg 命令
    ffmpeg_command = (
        f"ffmpeg -loglevel error -loop 1 -t {duration} -i {image_path}"
        f" -f lavfi -i anullsrc=cl=mono:r=24000 -shortest -y "
        f" -vf 'scale={size[0]}:{size[1]}' {output_file}"
        
    )
    
    # 运行 FFmpeg 命令
    subprocess.run(ffmpeg_command, shell=True, check=True)   
   
    
def sort_filelist(folder_path, extension):
    """
    对以下情形适用：
        1.mp4
        10.mp4
        2.mp4
        20.mp4

    """

    filenames_set = get_filenames_by_extension(folder_path, extension)
    
    # 排序
    sorted_filenames = sort_filenames_as_integers(filenames_set)
    
    # 组成文件完整路径
    files = add_path_to_filenames(folder_path, sorted_filenames)
    
    return files


# 获取文件夹下所有文件的文件名和扩展名，并保存到set中
def get_filenames_by_extension(folder_path, extension):
    filenames_set = set()
    
    # 遍历文件夹中的所有文件
    for root, dirs, files in os.walk(folder_path):
        for file in files:
            filename, ext = os.path.splitext(file)
            if ext == extension:
                filenames_set.add(filename)
    
    return filenames_set
    
    
# 将set里元素转换成整型后排序
def sort_filenames_as_integers(filenames_set):
    # 将文件名转换为整型并排序
    try:
        filenames_list = sorted(int(filename) for filename in filenames_set)
    except ValueError:
        raise ValueError("Some filenames are not integers and cannot be converted.")
    
    # 转换回字符串（如果需要保持字符串格式）
    sorted_filenames = [str(filename) for filename in filenames_list]
    
    return sorted_filenames
    
    
def add_path_to_filenames(folder_path, filenames):
    full_paths = []
    for filename in filenames:
        full_path = os.path.join(folder_path, filename + '.mp4')
        full_paths.append(full_path)
    return full_paths    
    
    
# 将列表中的每个元素重复 n 次
def repeat_elements(input_list, n):
    # 使用列表推导式将每个元素重复 n 次
    return [item for item in input_list for _ in range(n)]
    
    
# 将列表中的每个元素重复n次，并加上间隔视频
def add_elements_with_silence(input_list, silence_file):
    result = []
    for item in input_list:
        result.append(item)
        result.append(silence_file)
    
    return result


# 将多个mp3合成为一个
def concatenate_mp3_files(input_txt, output_file):
    
    # 构建 FFmpeg 命令
    ffmpeg_command = "ffmpeg -loglevel error -f concat -segment_time_metadata 1 -safe 0 -i " \
        + input_txt + \
        " -af aselect=concatdec_select,aresample=async=1 -y "\
        + output_file
    
    # 运行FFmpeg命令合并视频并覆盖已有的output.mp3
    subprocess.call(ffmpeg_command, shell=True)
    #print(f"Audios have been merged into {output_file}")
    
    
# 将多个短视频合成为一个视频
def concatenate_mp4_files(input_txt, output_file):
    
    # 构建 FFmpeg 命令
    ffmpeg_command = "ffmpeg -loglevel error -f concat -segment_time_metadata 1 -safe 0 -i " \
        + input_txt + \
        " -vf select=concatdec_select -af aselect=concatdec_select,aresample=async=1 -y "\
        + output_file
    
    # 运行FFmpeg命令合并视频并覆盖已有的output.mp4
    subprocess.call(ffmpeg_command, shell=True)
    print(f"Videos have been merged into {output_file}")    
            
        
# 删除单个文件
def delete_file(file_path):
    if os.path.exists(file_path):
        os.remove(file_path)
        print(f"{file_path} has been deleted")
    else:
        print(f"{file_path} does not exist")


# 删除文件夹及其内容
def delete_folder(folder_path):
    if os.path.exists(folder_path):
        shutil.rmtree(folder_path)
        print(f"{folder_path} and its contents have been deleted")
    else:
        print(f"{folder_path} does not exist")
   

# 创建空白静音视频
def create_silent_video(duration, size, output_file):
    # 创建一个白色的静音视频
    ffmpeg_command = (
        f"ffmpeg -loglevel error -f lavfi -i color=c=white:s={size[0]}x{size[1]}:d={duration} "
        f"-f lavfi -t {duration} -i anullsrc=r=24000:cl=mono "
        f" -c:v libx264 -c:a aac -y {output_file}"
    )
    # -i input.mp4：指定输入视频文件。
    # -f lavfi -t 3 -i anullsrc=r=44100:cl=stereo：
    # -f lavfi：指定输入格式为Libavfilter。
    # -t 3：将输入流的持续时间限制为3秒。
    # -i anullsrc=r=44100:cl=stereo：使用anullsrc滤镜生成一个静音音频流，采样率为44100 Hz，立体声。
    # -c:v copy：直接复制视频流，不重新编码。
    # -c:a aac：将音频流编码为AAC格式。
    subprocess.run(ffmpeg_command, shell=True, check=True)  
       

def create_silent_mp3(duration, output_file):
    # Construct the ffmpeg command
    command = [
        'ffmpeg', 
        '-loglevel', 'error',
        '-f', 'lavfi', 
        '-i', 'anullsrc=r=44100:cl=stereo', 
        '-t', str(duration), 
        '-y', # 覆盖原有文件
        output_file
    ]

    # Run the command
    subprocess.run(command, check=True)
    print("Done! Create silent mp3")
    return output_file
